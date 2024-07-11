import logging
import pathlib
import re
import tempfile
from typing import List, Tuple

import json5
import pptx
from global_config import GlobalConfig

PATTERN = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)
SAMPLE_JSON_FOR_PPTX = '''
{
    "title": "Understanding AI",
    "slides": [
        {
            "heading": "Introduction",
            "bullet_points": [
                "Brief overview of AI",
                [
                    "Importance of understanding AI"
                ]
            ]
        }
    ]
}
'''

logging.basicConfig(
    level=GlobalConfig.LOG_LEVEL,
    format='%(asctime)s - %(message)s',
)


def remove_slide_number_from_heading(header: str) -> str:
    """
    Remove the slide number from a given slide header.

    :param header: The header of a slide
    """

    if PATTERN.match(header):
        idx = header.find(':')
        header = header[idx + 1:]

    return header

def find_split_point(flat_items_list, max_chars_per_slide):
    """Find the best split point around the character threshold, preferably at level=1."""
    current_char_count = 0
    for i, (text, level) in enumerate(flat_items_list):
        current_char_count += len(text)
        if current_char_count > max_chars_per_slide:
            # Look for a level=1 item around the threshold
            for j in range(max(0, i-3), min(len(flat_items_list), i+3)):
                if flat_items_list[j][1] == 1:
                    return j
            return i
    return len(flat_items_list)

def split_slide_content(flat_items_list, max_chars_per_slide):
    """Split content into multiple slides if it exceeds the maximum characters per slide."""
    slides_content = []
    while flat_items_list:
        split_point = find_split_point(flat_items_list, max_chars_per_slide)
        slides_content.append(flat_items_list[:split_point])
        flat_items_list = flat_items_list[split_point:]
    return slides_content

def cleanup_slides_data(data_string):
    data_dict = None
    try:
        data_dict = json.loads(data_string)  # Convert JSON string to a dictionary
    except Exception as e:
        ee = f"Error parsing JSON data: {e} with data = {data_string}"
        print(ee)
        return data_string
    #
    if "slides" not in data_dict:
        return data_string
    #
        
    slides = data_dict["slides"]  # Access the slides array
    processed_slides = []  # Initialize an empty list to store the processed elements

    for element in slides:
        if isinstance(element, str):
            continue  # Skip if the element is a string
        elif isinstance(element, dict):
            processed_slides.append(element)  # Keep if the element is an object
        elif isinstance(element, list):
            for item in element:
                if isinstance(item, dict):
                    processed_slides.append(item)  # Insert each object from the array into the processed list

    data_dict["slides"] = processed_slides  # Update the slides array with the processed elements
    return json.dumps(data_dict, indent=4)  # Convert dictionary back to JSON string and return it
#


def generate_powerpoint_presentation(structured_data: str, slides_template: str, output_file_path: pathlib.Path, max_chars_per_slide: int = 1000) -> List:
    """
    Create and save a PowerPoint presentation file containing the content in JSON format.

    :param structured_data: The presentation contents as "JSON" (may contain trailing commas)
    :param slides_template: The PPTX template to use
    :param output_file_path: The path of the PPTX file to save as
    :param max_chars_per_slide: Maximum number of characters allowed per slide before splitting
    :return A list of presentation title and slides headers
    """

    BAD_STRING = "This is an example response from ChatGPT.,"
    structured_data = structured_data.replace(BAD_STRING, "")
    structured_data = cleanup_slides_data(structured_data)

    # The structured "JSON" might contain trailing commas, so using json5
    parsed_data = None
    try:
        parsed_data = json5.loads(structured_data)
    except Exception as e:
        ee = f"Error parsing JSON5 data: {e} with data = {structured_data}"
        logging.error(ee)
        raise ValueError(ee)
    #
            
    config = GlobalConfig()
    
    logging.debug(
        "*** Using PPTX template: %s",
        config.PPTX_TEMPLATE_FILES[slides_template]['file']
    )
    presentation = pptx.Presentation(config.PPTX_TEMPLATE_FILES[slides_template]['file'])

    # The title slide
    if len(presentation.slides) > 0:
        for i in range(len(presentation.slides)-1, -1, -1): 
            rId = presentation.slides._sldIdLst[i].rId
            presentation.part.drop_rel(rId)
            del presentation.slides._sldIdLst[i]
        #
    #
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = parsed_data['title']
    logging.debug('Presentation title is: %s', title.text)
    if 'subtitle' in parsed_data:
        subtitle.text = parsed_data['subtitle']
        logging.debug('Presentation subtitle is: %s', subtitle.text)
    else:
        subtitle.text = ''
    #
    all_headers = [title.text, ]

    slides_parsed_data = parsed_data.get('slides', None)
    if slides_parsed_data is None:
        error_message = 'No slides found in the parsed data'
        logging.error('error_message')
        print(error_message)
        return []
    #



    
    new_slides_data = []

    for a_slide in slides_parsed_data:
        flat_items_list = get_flat_list_of_contents(a_slide['bullet_points'], level=0)
        slides_content = split_slide_content(flat_items_list, max_chars_per_slide)
        
        for i, slide_content in enumerate(slides_content):
            if i == 0:
                slide_heading = a_slide['heading']
            else:
                slide_heading = f"{a_slide['heading']} (continued {i})" if i > 1 else f"{a_slide['heading']} (continued)"
            slide_data = {
                'heading': slide_heading,
                'bullet_points': slide_content
            }
            new_slides_data.append(slide_data)

    for slide_data in new_slides_data:
        if "type" in slide_data and slide_data["type"] == "sectionheader":
            bullet_slide_layout = presentation.slide_layouts[2]
        else:
            bullet_slide_layout = presentation.slide_layouts[1]
        
        slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = remove_slide_number_from_heading(slide_data['heading'])
        all_headers.append(title_shape.text)
        text_frame = body_shape.text_frame

        for text, level in slide_data['bullet_points']:
            paragraph = text_frame.add_paragraph()
            
            while text:
                hyperlink_start = text.find('https://') if 'https://' in text else text.find('http://')
                if hyperlink_start == -1:
                    paragraph.add_run().text = text
                    break

                # Add text before hyperlink
                if hyperlink_start > 0:
                    paragraph.add_run().text = text[:hyperlink_start]
                    text = text[hyperlink_start:]

                # Add hyperlink
                hyperlink_end = text.find(' ', hyperlink_start)
                if hyperlink_end == -1:
                    hyperlink_url = text
                    text = ''
                else:
                    hyperlink_url = text[:hyperlink_end]
                    text = text[hyperlink_end + 1:]

                run = paragraph.add_run()
                run.text = hyperlink_url
                hlink = run.hyperlink
                hlink.address = hyperlink_url

            paragraph.level = level

    presentation.save(output_file_path)

    return all_headers
def generate_powerpoint_presentation_advanced(
        structured_data: str,
        slides_template: str,
        output_file_path: pathlib.Path
) -> List:
    """
    Create and save a PowerPoint presentation file containing the content in JSON format.

    :param structured_data: The presentation contents as "JSON" (may contain trailing commas)
    :param slides_template: The PPTX template to use
    :param output_file_path: The path of the PPTX file to save as
    :return A list of presentation title and slides headers
    """

    # The structured "JSON" might contain trailing commas, so using json5
    parsed_data = None
    try:
        parsed_data = json5.loads(structured_data)
    except Exception as e:
        print("Error parsing JSON5 data: ", e)
        print("structured_data = ", structured_data)
        return []
    #
    print("--------")
    print("parsed_data = ", parsed_data)
    print("--------")
    
    data_slides = []
    data_title = ""
    data_subtitle = ""
    if 'slides' in parsed_data: 
        data_slides = parsed_data['slides'] 
    else: 
        print(f"[ERROR] No slides in parsed_data: f{parsed_data}")
    #
    
    if 'title' in parsed_data: data_title = parsed_data['title']
    if 'subtitle' in parsed_data: data_subtitle = parsed_data['subtitle']
    
    print("==============")
    print("data_slides = ", data_slides)
    print("==============")
    print(f"data_title = {data_title}, data_subtitle = {data_subtitle}")
    print("==============")
      
    slides = json5.loads(data_slides)  
    config = GlobalConfig()
    
    logging.debug(
        "*** Using PPTX template: %s",
        config.PPTX_TEMPLATE_FILES[slides_template]['file']
    )
    presentation = pptx.Presentation(config.PPTX_TEMPLATE_FILES[slides_template]['file'])

    # The title slide
    pptx_slide = None
    if len(presentation.slides) > 0:
        for i in range(len(presentation.slides)-1, -1, -1): 
            rId = presentation.slides._sldIdLst[i].rId
            presentation.part.drop_rel(rId)
            del presentation.slides._sldIdLst[i]
        #
    #
    title_slide_layout = presentation.slide_layouts[0]
    pptx_slide = presentation.slides.add_slide(title_slide_layout)
    pptx_title = pptx_slide.shapes.title
    pptx_title.text = data_title
    pptx_subtitle = pptx_slide.placeholders[1]
    pptx_subtitle.text = data_subtitle

    all_headers = [pptx_title.text, ]
    print("all_headers = ", all_headers)

    # Add contents in a loop   
    lettered = ["A", "B", "C", "D", "E", "F", "G", "H", "I", "J", "K", "L", "M", "N", "O", "P", "Q", "R", "S", "T", "U", "V", "W", "X", "Y", "Z"]
    for a_slide in slides:
        numbered_indexes = [0,0,0,0,0]
        lettered_indexes = [0,0,0,0,0]
        print("------------------")
        print(f"slide = {a_slide}")
        
        heading = ""
        bullet_points = []
        section_type = ""
        is_section_header = False
        
        if "heading" in a_slide: heading = a_slide['heading']
        if "bullet_points" in a_slide: bullet_points = a_slide['bullet_points']
        if "type" in a_slide: section_type = a_slide['type']
        
        if section_type == "sectionheader": 
            is_section_header = True
            print("This is a section header")
            bullet_slide_layout = presentation.slide_layouts[2]
        else:
            bullet_slide_layout = presentation.slide_layouts[1]
        #
        
        current_slide = presentation.slides.add_slide(bullet_slide_layout)
        current_shapes = current_slide.shapes
        title_shape = current_shapes.title            
        body_shape = current_shapes.placeholders[1]
        title_shape.text = remove_slide_number_from_heading(heading)
        all_headers.append(title_shape.text)
        text_frame = body_shape.text_frame

        if len(bullet_points) > 0:
            
            
            # The bullet_points is a flat array with indentation level explicitly set
            # We create bullets, numbers and letters 'bullets' manually
            previous_numbered_level = -1
            previous_lettered_level = -1
            
            for an_item in bullet_points:
                
                bullet_type = "none"
                bullet_level = 0
                bullet_text = ""
                
                if "bullet_type" in an_item: bullet_type = an_item["bullet_type"]    
                if "bullet_level" in an_item: bullet_level = int(an_item["bullet_level"])
                if "bullet_text" in an_item: bullet_text = an_item["bullet_text"]
                
                indentation = bullet_level
                if bullet_type == "bullet":
                    bullet_text = "• " + bullet_text
                elif bullet_type == "number":
                    if previous_numbered_level < bullet_level: numbered_indexes[bullet_level] = 0
                    bullet_text = f"{numbered_indexes[bullet_level] + 1}. " + bullet_text
                    previous_numbered_level = bullet_level
                    numbered_indexes[bullet_level] += 1
                    #indentation += 2
                elif bullet_type == "letter":
                    if previous_lettered_level < bullet_level: lettered_indexes[bullet_level] = 0
                    bullet_text = f"{lettered[lettered_indexes[bullet_level]]}. " + bullet_text
                    lettered_indexes[bullet_level] += 1
                    previous_lettered_level = bullet_level
                    #indentation += 3
                #        
                
                print(f"bullet_type = {bullet_type}, bullet_level = {bullet_level}, bullet_text = {bullet_text}, indentation = {indentation}")           
                paragraph = text_frame.add_paragraph()
                paragraph.text = bullet_text
                paragraph.level = indentation
            #
        else:
            # create a blank bullet point
            paragraph = text_frame.add_paragraph()
            paragraph.text = ""
            paragraph.level = 0
        #
    #

    presentation.save(output_file_path)

    return all_headers
#



def get_flat_list_of_contents(items: list, level: int) -> List[Tuple]:
    """
    Flatten a (hierarchical) list of bullet points to a single list containing each item and
    its level.

    :param items: A bullet point (string or list)
    :param level: The current level of hierarchy
    :return: A list of (bullet item text, hierarchical level) tuples
    """

    flat_list = []

    for item in items:
        if isinstance(item, str):
            flat_list.append((item, level))
        elif isinstance(item, list):
            flat_list = flat_list + get_flat_list_of_contents(item, level + 1)

    return flat_list


if __name__ == '__main__':
    # bullets = [
    #     'Description',
    #     'Types',
    #     [
    #         'Type A',
    #         'Type B'
    #     ],
    #     'Grand parent',
    #     [
    #         'Parent',
    #         [
    #             'Grand child'
    #         ]
    #     ]
    # ]

    # output = get_flat_list_of_contents(bullets, level=0)
    # for x in output:
    #     print(x)

    json_data = '''
    {
    "title": "Understanding AI",
    "slides": [
        {
            "heading": "Introduction",
            "bullet_points": [
                "Brief overview of AI",
                [
                    "Importance of understanding AI"
                ]
            ]
        },
        {
            "heading": "What is AI?",
            "bullet_points": [
                "Definition of AI",
                [
                    "Types of AI",
                    [
                        "Narrow or weak AI",
                        "General or strong AI"
                    ]
                ],
                "Differences between AI and machine learning"
            ]
        },
        {
            "heading": "How AI Works",
            "bullet_points": [
                "Overview of AI algorithms",
                [
                    "Types of AI algorithms",
                    [
                        "Rule-based systems",
                        "Decision tree systems",
                        "Neural networks"
                    ]
                ],
                "How AI processes data"
            ]
        },
        {
            "heading": "Pros of AI",
            "bullet_points": [
                "Increased efficiency and productivity",
                "Improved accuracy and precision",
                "Enhanced decision-making capabilities",
                "Personalized experiences"
            ]
        },
        {
            "heading": "Cons of AI",
            "bullet_points": [
                "Job displacement and loss of employment",
                "Bias and discrimination",
                "Privacy and security concerns",
                "Dependence on technology"
            ]
        },
        {
            "heading": "Future Prospects of AI",
            "bullet_points": [
                "Advancements in fields such as healthcare and finance",
                "Increased use"
            ]
        }
    ]
}'''

    temp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    path = pathlib.Path(temp.name)

    generate_powerpoint_presentation(
        json5.loads(json_data),
        output_file_path=path,
        slides_template='Blank'
    )
